"""File handlers for different document types."""

import io
from pathlib import Path
from typing import Optional, Dict, Any
from abc import ABC, abstractmethod

import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd


class FileHandler(ABC):
    """Abstract base class for file handlers."""
    
    @abstractmethod
    def extract_text(self, file_path: Path) -> str:
        """Extract text from file."""
        pass
    
    @abstractmethod
    def get_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from file."""
        pass


class PDFHandler(FileHandler):
    """Handler for PDF files."""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise ValueError(f"Error extracting text from PDF: {str(e)}")
        
        return text.strip()
    
    def get_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from PDF file."""
        metadata = {}
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                if pdf_reader.metadata:
                    metadata = {
                        "title": pdf_reader.metadata.get('/Title', ''),
                        "author": pdf_reader.metadata.get('/Author', ''),
                        "subject": pdf_reader.metadata.get('/Subject', ''),
                        "creator": pdf_reader.metadata.get('/Creator', ''),
                        "pages": len(pdf_reader.pages)
                    }
        except Exception:
            pass
        
        return metadata


class DocxHandler(FileHandler):
    """Handler for DOCX files."""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        text = ""
        try:
            doc = DocxDocument(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\t"
                    text += "\n"
        except Exception as e:
            raise ValueError(f"Error extracting text from DOCX: {str(e)}")
        
        return text.strip()
    
    def get_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from DOCX file."""
        metadata = {}
        try:
            doc = DocxDocument(file_path)
            core_props = doc.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else ""
            }
        except Exception:
            pass
        
        return metadata


class PptxHandler(FileHandler):
    """Handler for PPTX files."""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from PPTX file."""
        text = ""
        try:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise ValueError(f"Error extracting text from PPTX: {str(e)}")
        
        return text.strip()
    
    def get_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from PPTX file."""
        metadata = {}
        try:
            presentation = Presentation(file_path)
            core_props = presentation.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": str(core_props.created) if core_props.created else "",
                "modified": str(core_props.modified) if core_props.modified else "",
                "slides": len(presentation.slides)
            }
        except Exception:
            pass
        
        return metadata


class TxtHandler(FileHandler):
    """Handler for TXT files."""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
        except UnicodeDecodeError:
            # Try with different encoding if UTF-8 fails
            with open(file_path, 'r', encoding='latin-1') as file:
                text = file.read()
        except Exception as e:
            raise ValueError(f"Error extracting text from TXT: {str(e)}")
        
        return text.strip()
    
    def get_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from TXT file."""
        stats = file_path.stat()
        return {
            "size": stats.st_size,
            "modified": stats.st_mtime,
            "created": stats.st_ctime
        }


class XlsxHandler(FileHandler):
    """Handler for XLSX files."""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from XLSX file."""
        text = ""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += df.to_string(index=False) + "\n"
        except Exception as e:
            raise ValueError(f"Error extracting text from XLSX: {str(e)}")
        
        return text.strip()
    
    def get_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from XLSX file."""
        metadata = {}
        try:
            excel_file = pd.ExcelFile(file_path)
            metadata = {
                "sheets": excel_file.sheet_names,
                "num_sheets": len(excel_file.sheet_names)
            }
        except Exception:
            pass
        
        return metadata


class FileHandlerFactory:
    """Factory for creating appropriate file handlers."""
    
    _handlers = {
        ".pdf": PDFHandler,
        ".docx": DocxHandler,
        ".pptx": PptxHandler,
        ".txt": TxtHandler,
        ".xlsx": XlsxHandler,
        ".xls": XlsxHandler
    }
    
    @classmethod
    def get_handler(cls, file_path: Path) -> FileHandler:
        """Get appropriate handler for file type."""
        extension = file_path.suffix.lower()
        handler_class = cls._handlers.get(extension)
        
        if not handler_class:
            raise ValueError(f"Unsupported file type: {extension}")
        
        return handler_class()